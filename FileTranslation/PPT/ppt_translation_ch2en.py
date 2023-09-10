from pptx import Presentation
from time import sleep
import transformers
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import os
import sys
if __name__ == '__main__':
    savePath = sys.argv[2]
    # 1、ppt地址
    pptPath = sys.argv[1]
    model_path = 'models/zh-en'
    tokenizer = AutoTokenizer.from_pretrained(model_path)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_path)
    translate_model = transformers.pipeline("translation", model=model, tokenizer=tokenizer)
    prs = Presentation(pptPath)

    print("start")
    for ns, slide in enumerate(prs.slides):
        for nsh, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for np, paragraph in enumerate(shape.text_frame.paragraphs):
                for rs, run in enumerate(paragraph.runs):
                    str_in = run.text
                    str_out = translate_model(str_in)[0]['translation_text']
                    prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = str_out
                    sleep(1.5)
                    print(np)
    ppt_name = os.path.basename(pptPath)
    print(ppt_name)
    prs.save(os.path.join(savePath,ppt_name))
    print(os.path.join(savePath,ppt_name))