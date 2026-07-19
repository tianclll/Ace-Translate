# Copyright (c) 2025 PaddlePaddle Authors. All Rights Reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
from pathlib import Path

from .base import BaseConverter, ConvertResult
from .math import (
    extract_math_from_paragraph as _extract_math_from_paragraph,
    paragraph_has_math as _paragraph_has_math,
)
from .registry import default_registry


# pptx XML namespace for DrawingML run properties
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A = "{" + _A_NS + "}"
_A_STRIKE = f"{_A}strike"
_A_RPR = f"{_A}rPr"
_A_P = f"{_A}p"

# Markup Compatibility namespace
_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_MC = "{" + _MC_NS + "}"
_MC_CHOICE = f"{_MC}Choice"
_MC_ALT = f"{_MC}AlternateContent"

# Chart namespace
_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_C = "{" + _C_NS + "}"

# Chart type name lookup (used in _chart_to_html)
_CHART_TYPE_NAMES = {
    1: "Area Chart",  # AREA
    2: "Area Chart",  # AREA_STACKED
    4: "Line Chart",  # LINE
    5: "Pie Chart",  # PIE
    15: "Bubble Chart",  # BUBBLE
    51: "Column Chart",  # COLUMN_CLUSTERED
    52: "Column Chart",  # COLUMN_STACKED
    53: "Column Chart",  # COLUMN_STACKED_100
    57: "Bar Chart",  # BAR_CLUSTERED
    58: "Bar Chart",  # BAR_STACKED
    65: "Line Chart",  # LINE_MARKERS
    -4120: "Doughnut Chart",  # DOUGHNUT
    -4151: "Radar Chart",  # RADAR
    -4169: "Scatter Chart",  # XY_SCATTER
}


def _pptx_run_strike(run) -> bool:
    """Return True if the run has strikethrough (sngStrike or dblStrike) set in XML."""
    try:
        rPr = run._r.find(_A_RPR)
        if rPr is not None:
            val = rPr.get(_A_STRIKE)
            return val in ("sngStrike", "dblStrike")
    except Exception:
        pass
    return False


def _pptx_run_script(run) -> str:
    """Return 'super', 'sub', or '' based on DrawingML baseline attribute.

    baseline > 0  -> superscript, baseline < 0 -> subscript.
    baseline attribute is on <a:rPr> with no namespace prefix.
    """
    try:
        rPr = run._r.find(_A_RPR)
        if rPr is not None:
            val = rPr.get("baseline")
            if val is not None:
                n = int(val)
                if n > 0:
                    return "super"
                if n < 0:
                    return "sub"
    except Exception:
        pass
    return ""


def _escape_md_url(url: str) -> str:
    """Escape parentheses in URL for Markdown link syntax."""
    return url.replace("(", "%28").replace(")", "%29")


def _classify_part(part: str) -> str:
    """Classify a slide content part for grouping (heading/html/list/blockquote/other)."""
    s = part.lstrip()
    if s.startswith("##"):
        return "heading"
    if s.startswith("<"):
        return "html"
    if s.startswith("- "):
        return "list"
    if s.startswith(">"):
        return "blockquote"
    return "other"


def _format_run_segment(
    seg: str,
    bold: bool,
    italic: bool,
    underline: bool,
    strikethrough: bool,
    script: str,
    url: str,
) -> str:
    """Apply Markdown/HTML inline formatting to a text segment."""
    t = seg
    if bold or italic or underline or strikethrough or script:
        leading = len(t) - len(t.lstrip())
        trailing = len(t) - len(t.rstrip())
        prefix = t[:leading] if leading else ""
        suffix = t[len(t) - trailing :] if trailing else ""
        inner = t.strip()
        if inner:
            if strikethrough:
                inner = f"~~{inner}~~"
            if bold and italic:
                inner = f"***{inner}***"
            elif bold:
                inner = f"**{inner}**"
            elif italic:
                inner = f"*{inner}*"
            if underline:
                inner = f"<u>{inner}</u>"
            if script == "super":
                inner = f"<sup>{inner}</sup>"
            elif script == "sub":
                inner = f"<sub>{inner}</sub>"
            t = prefix + inner + suffix
        elif underline and t:
            # Pure whitespace + underline = fill-in line
            # Replace spaces with NBSP so Markdown renderers preserve width
            t = "<u>" + "\u00a0" * len(t) + "</u>"
    if url:
        escaped_url = _escape_md_url(url)
        t = f"[{t}]({escaped_url})"
    return t


@default_registry.register
class PptxConverter(BaseConverter):
    supported_extensions = ["pptx"]
    supported_mimetypes = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    def convert_file(self, file_path: Path, **kwargs) -> ConvertResult:
        try:
            from pptx import Presentation
            from pptx.shapes.picture import Picture
        except ImportError:
            raise RuntimeError(
                "PPTX conversion requires python-pptx: pip install paddleocr[doc2md]"
            )
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            MSO_SHAPE_TYPE = None  # type: ignore[assignment]

        prs = Presentation(str(file_path))
        slides_md = []
        images: dict = {}
        image_counter = [0]
        slide_width = prs.slide_width
        self._Picture = Picture
        self._MSO_SHAPE_TYPE = MSO_SHAPE_TYPE

        for slide in prs.slides:
            slide_parts = []

            # Process all shapes
            for shape in slide.shapes:
                self._process_shape(
                    shape, slide_parts, images, image_counter, slide_width, slide.part
                )

            # Handle math formulas inside mc:AlternateContent elements
            # (python-pptx doesn't expose these as Shape objects)
            for alt_content in slide._element.iter(_MC_ALT):
                # Only look at mc:Choice (the preferred rendering path)
                choice = alt_content.find(_MC_CHOICE)
                if choice is None:
                    continue
                for para_elem in choice.iter(_A_P):
                    if _paragraph_has_math(para_elem):
                        math_items = _extract_math_from_paragraph(para_elem)
                        for latex in math_items:
                            slide_parts.append(f"$$\n{latex}\n$$")

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **Notes**: {notes_text}")

            # Group parts by content type and separate groups with blank lines
            # to prevent HTML blocks from consuming adjacent list items
            groups: list[list[str]] = []
            for part in slide_parts:
                kind = _classify_part(part)
                if groups and _classify_part(groups[-1][0]) == kind:
                    groups[-1].append(part)
                else:
                    groups.append([part])

            slides_md.append("\n\n".join("\n".join(g) for g in groups))

        md_text = "\n\n---\n\n".join(slides_md)

        return ConvertResult(
            markdown=md_text,
            images=images,
            metadata={
                "format": "PPTX",
                "slide_count": len(prs.slides),
            },
        )

    def _process_shape(
        self, shape, slide_parts, images, image_counter, slide_width, slide_part
    ):
        """Recursively process a shape: Picture, GroupShape, Chart, Table, or TextFrame."""
        Picture = self._Picture
        MSO_SHAPE_TYPE = self._MSO_SHAPE_TYPE

        # 1. Picture
        if isinstance(shape, Picture):
            try:
                img = shape.image
                image_counter[0] += 1
                filename = f"image{image_counter[0]}.{img.ext}"
                rel_path = f"assets/{filename}"
                images[rel_path] = img.blob
                if shape.width and slide_width:
                    pct = min(round(shape.width / slide_width * 100), 100)
                    slide_parts.append(f'<img src="assets/{filename}" width="{pct}%">')
                else:
                    slide_parts.append(f'<img src="assets/{filename}">')
            except (ValueError, AttributeError):
                pass
            return

        # 2. GroupShape - recurse into child shapes
        if MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    self._process_shape(
                        child,
                        slide_parts,
                        images,
                        image_counter,
                        slide_width,
                        slide_part,
                    )
            except AttributeError:
                pass
            return

        # 3. Chart
        if shape.has_chart:
            slide_parts.append(self._chart_to_html(shape.chart))
            return

        # 4. Table
        if shape.has_table:
            slide_parts.append(
                self._table_to_html(shape.table, slide_part, image_counter, images)
            )
            return

        # 5. TextFrame
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # Check for math elements first
                para_xml = paragraph._p
                if _paragraph_has_math(para_xml):
                    math_items = _extract_math_from_paragraph(para_xml)
                    for latex in math_items:
                        slide_parts.append(f"$$\n{latex}\n$$")
                    continue

                parts = []
                for run in paragraph.runs:
                    t = run.text
                    if not t:
                        continue
                    try:
                        url = run.hyperlink.address
                    except Exception:
                        url = None

                    bold = bool(run.font.bold)
                    italic = bool(run.font.italic)
                    underline = bool(run.font.underline) and not url
                    strikethrough = _pptx_run_strike(run)
                    script = _pptx_run_script(run)

                    if "\n" in t:
                        segments = t.split("\n")
                        for j, seg in enumerate(segments):
                            if seg:
                                parts.append(
                                    _format_run_segment(
                                        seg,
                                        bold,
                                        italic,
                                        underline,
                                        strikethrough,
                                        script,
                                        url,
                                    )
                                )
                            if j < len(segments) - 1:
                                parts.append("<br>\n")
                    else:
                        parts.append(
                            _format_run_segment(
                                t, bold, italic, underline, strikethrough, script, url
                            )
                        )
                text = "".join(parts).strip()
                if not text:
                    continue
                level = paragraph.level
                indent = "  " * level
                # List item continuation lines need indentation; revert <br>\n → <br>
                text = text.replace("<br>\n", "<br>")
                slide_parts.append(f"{indent}- {text}")

    def _chart_to_html(self, chart) -> str:
        """Extract chart data as an HTML table."""
        try:
            chart_type_val = chart.chart_type.value if chart.chart_type else 0
            chart_type_name = _CHART_TYPE_NAMES.get(chart_type_val, "Chart")
        except Exception:
            chart_type_name = "Chart"

        try:
            title_text = ""
            try:
                title_text = chart.chart_title.text_frame.text.strip()
            except Exception:
                pass

            # Extract axis info from OOXML
            chart_root = chart._element

            cat_ax_title = ""
            for ax_tag in (f"{_C}catAx", f"{_C}dateAx"):
                ax_el = chart_root.find(f".//{ax_tag}")
                if ax_el is not None:
                    t_el = ax_el.find(f"{_C}title")
                    if t_el is not None:
                        texts = t_el.findall(f".//{_A}t")
                        cat_ax_title = "".join(el.text or "" for el in texts).strip()
                    break

            has_date_ax = chart_root.find(f".//{_C}dateAx") is not None

            val_ax_title = ""
            val_ax_el = chart_root.find(f".//{_C}valAx")
            if val_ax_el is not None:
                t_el = val_ax_el.find(f"{_C}title")
                if t_el is not None:
                    texts = t_el.findall(f".//{_A}t")
                    val_ax_title = "".join(el.text or "" for el in texts).strip()

            plot = chart.plots[0]
            categories = list(plot.categories) if plot.categories else []
            series_list = list(plot.series)

            if not series_list:
                return f"[{chart_type_name}]"

            # Convert Excel date serials to YYYY-MM-DD for date axes
            if has_date_ax and categories:
                from datetime import datetime, timedelta

                converted = []
                for c in categories:
                    try:
                        dt = datetime(1899, 12, 30) + timedelta(days=float(c))
                        converted.append(dt.strftime("%Y-%m-%d"))
                    except (ValueError, TypeError):
                        converted.append(str(c) if c is not None else "")
                categories = converted

            # Collect series names and values
            series_names = []
            series_values = []
            for idx, series in enumerate(series_list):
                try:
                    name = (
                        (series.tx.text if series.tx else "")
                        or val_ax_title
                        or f"Series{idx+1}"
                    )
                except Exception:
                    name = val_ax_title or f"Series{idx+1}"
                series_names.append(name)
                try:
                    vals = [
                        str(round(v, 4)) if v is not None else "" for v in series.values
                    ]
                except Exception:
                    vals = []
                series_values.append(vals)

            # Build HTML table
            html_parts = ["<table>"]
            if title_text:
                html_parts.append(f"<caption>{title_text}</caption>")

            has_header = cat_ax_title or any(name for name in series_names)
            if has_header:
                html_parts.append("<thead><tr>")
                html_parts.append(f"<th>{cat_ax_title}</th>")
                for name in series_names:
                    html_parts.append(f"<th>{name}</th>")
                html_parts.append("</tr></thead>")

            html_parts.append("<tbody>")
            if categories:
                for i, cat in enumerate(categories):
                    html_parts.append(f"<tr><td>{cat}</td>")
                    for vals in series_values:
                        v = vals[i] if i < len(vals) else ""
                        html_parts.append(f"<td>{v}</td>")
                    html_parts.append("</tr>")
            else:
                max_len = max((len(v) for v in series_values), default=0)
                for i in range(max_len):
                    html_parts.append(f"<tr><td>Item{i+1}</td>")
                    for vals in series_values:
                        v = vals[i] if i < len(vals) else ""
                        html_parts.append(f"<td>{v}</td>")
                    html_parts.append("</tr>")
            html_parts.append("</tbody></table>")

            return "\n".join(html_parts)
        except Exception:
            return f"[{chart_type_name}]"

    @staticmethod
    def _table_to_html(
        table, slide_part, image_counter_list: list, images: dict
    ) -> str:
        """Convert a PPTX table to an HTML table, handling merged cells and cell background images."""
        _BLIP_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        _REL_NS = (
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        )

        visited: set[tuple[int, int]] = set()
        html_parts = ["<table>"]

        for i, row in enumerate(table.rows):
            html_parts.append("<tr>")
            for j, cell in enumerate(row.cells):
                if (i, j) in visited:
                    continue
                tag = "th" if i == 0 else "td"
                attrs = ""
                if cell.is_merge_origin:
                    rs = cell.span_height
                    cs = cell.span_width
                    if cs > 1:
                        attrs += f' colspan="{cs}"'
                    if rs > 1:
                        attrs += f' rowspan="{rs}"'
                    for di in range(rs):
                        for dj in range(cs):
                            if (di, dj) != (0, 0):
                                visited.add((i + di, j + dj))

                content_parts = []

                # Extract cell background blip images
                blips = cell._tc.findall(f".//{_BLIP_NS}blip")
                for blip in blips:
                    r_embed = blip.get(f"{_REL_NS}embed")
                    if r_embed:
                        try:
                            image_part = slide_part.related_parts[r_embed]
                            image_counter_list[0] += 1
                            ext = image_part.content_type.split("/")[-1]
                            filename = f"image{image_counter_list[0]}.{ext}"
                            rel_path = f"assets/{filename}"
                            images[rel_path] = image_part.blob
                            content_parts.append(
                                f'<img src="assets/{filename}" width="100%">'
                            )
                        except (KeyError, AttributeError):
                            pass

                cell_text_parts = []
                for para in cell.text_frame.paragraphs:
                    # Check for math elements first
                    para_xml = para._p
                    if _paragraph_has_math(para_xml):
                        math_items = _extract_math_from_paragraph(para_xml)
                        for latex in math_items:
                            cell_text_parts.append(f"${latex}$")
                        continue

                    run_parts = []
                    for run in para.runs:
                        t = run.text or ""
                        if not t:
                            continue
                        try:
                            url = run.hyperlink.address
                        except Exception:
                            url = None

                        bold = bool(run.font.bold)
                        italic = bool(run.font.italic)
                        underline = bool(run.font.underline) and not url
                        strikethrough = bool(run.font.strike)
                        script = _pptx_run_script(run)

                        if bold:
                            t = f"<b>{t}</b>"
                        if italic:
                            t = f"<i>{t}</i>"
                        if underline:
                            t = f"<u>{t}</u>"
                        if strikethrough:
                            t = f"<del>{t}</del>"
                        if script == "super":
                            t = f"<sup>{t}</sup>"
                        elif script == "sub":
                            t = f"<sub>{t}</sub>"

                        if url:
                            run_parts.append(f'<a href="{url}">{t}</a>')
                        else:
                            run_parts.append(t)
                    cell_text_parts.append("".join(run_parts))
                text = "<br>".join(p for p in cell_text_parts if p.strip())
                if text:
                    content_parts.append(text)
                cell_html = "<br>".join(content_parts) if content_parts else ""

                html_parts.append(f"<{tag}{attrs}>{cell_html}</{tag}>")
            html_parts.append("</tr>")

        html_parts.append("</table>")
        return "\n".join(html_parts)
